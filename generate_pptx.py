from pptx import Presentation
from pptx.util import Pt
import re

def create_presentation(md_file, pptx_file):
    with open(md_file, 'r', encoding='utf-8') as f:
        lines = f.readlines()

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # Main Title
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "2026 AI AGENT 융합·확산 지원사업\nKick-off 미팅 자료"
    slide.placeholders[1].text = "수행계획서 작성 R&R 및 핵심 차별화 전략"

    current_slide = None
    tf = None

    for line in lines:
        stripped_line = line.strip()
        if not stripped_line:
            continue
            
        if line.startswith('## '):
            # Section Title Slide
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = stripped_line[3:]
            current_slide = None
            
        elif line.startswith('### '):
            # Content Slide
            slide = prs.slides.add_slide(bullet_slide_layout)
            slide.shapes.title.text = stripped_line[4:]
            tf = slide.placeholders[1].text_frame
            tf.clear()
            current_slide = slide
            
        elif stripped_line.startswith('- '):
            if current_slide:
                leading_spaces = len(line) - len(line.lstrip())
                level = 0
                if leading_spaces >= 2:
                    level = 1
                if leading_spaces >= 4:
                    level = 2
                    
                content = stripped_line[2:].replace('**', '')
                p = tf.add_paragraph()
                p.text = content
                p.level = level
                
        elif re.match(r'^\d\)', stripped_line):
            if current_slide:
                content = stripped_line.replace('**', '')
                p = tf.add_paragraph()
                p.text = content
                p.level = 1

    prs.save(pptx_file)

if __name__ == '__main__':
    create_presentation('Kickoff_발표자료_본문.md', 'Kickoff_발표자료_생성본.pptx')
